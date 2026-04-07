"""
AEF HRM – 6-Slide PowerPoint (5-min presentation + 5-min live demo)
Run: python build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

CYAN  = RGBColor(0x31, 0xb8, 0xcf)
NAVY  = RGBColor(0x0A, 0x4D, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1a, 0x1a, 0x2e)
LGRAY = RGBColor(0xf0, 0xf4, 0xf8)
MGRAY = RGBColor(0x71, 0x8E, 0xA4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

LOGO = os.path.join(os.path.dirname(__file__), 'static', 'LOGO.png')
LOGO_EXISTS = os.path.exists(LOGO)


def box(slide, l, t, w, h, fill, line=False):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = fill
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, l, t, w, h, size=16, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color


def logo(slide, l=0.3, t=0.15, w=1.7, h=0.62):
    if LOGO_EXISTS:
        slide.shapes.add_picture(LOGO, Inches(l), Inches(t), Inches(w), Inches(h))


def top_bar(slide):
    box(slide, 0, 0, 13.33, 7.5, LGRAY)
    box(slide, 0, 0, 13.33, 0.07, CYAN)
    logo(slide)


def bottom_bar(slide):
    box(slide, 0, 7.28, 13.33, 0.22, NAVY)
    txt(slide, 'Africa Eye Foundation  ·  AEF HR Management System  ·  Confidential',
        0.2, 7.28, 12.9, 0.22, size=9, color=WHITE, align=PP_ALIGN.CENTER)


def title_bar(slide, title, subtitle=None):
    box(slide, 0, 0.07, 13.33, 1.0, NAVY)
    txt(slide, title, 2.2, 0.1, 10.5, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 2.2, 0.78, 10.5, 0.32, size=13, color=CYAN, italic=True)


def divider(slide, t):
    d = slide.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(12.5), Inches(0.018))
    d.fill.solid(); d.fill.fore_color.rgb = CYAN; d.line.fill.background()


def bullet(slide, items, l=0.6, t=1.3, size=15, gap=0.54):
    for i, item in enumerate(items):
        txt(slide, f'›  {item}', l, t + i * gap, 12.0, gap, size=size, color=BLACK)


def card(slide, l, t, w, h, head, body, hbg=NAVY, bbg=WHITE):
    box(slide, l, t, w, 0.42, hbg)
    txt(slide, head, l + 0.1, t, w - 0.15, 0.42, size=12, bold=True, color=WHITE)
    box(slide, l, t + 0.42, w, h - 0.42, bbg)
    txt(slide, body, l + 0.12, t + 0.48, w - 0.22, h - 0.52, size=12, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.33, 7.5, NAVY)
box(s, 0, 0, 13.33, 1.3, CYAN)
box(s, 0, 6.2, 13.33, 1.3, CYAN)
logo(s, l=0.4, t=0.28, w=2.1, h=0.76)

txt(s, 'AEF HR MANAGEMENT SYSTEM',
    0.5, 1.65, 12.3, 1.0, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 'LeaveDesk — Integrated Human Resource Platform',
    0.5, 2.7, 12.3, 0.55, size=21, color=CYAN, align=PP_ALIGN.CENTER, italic=True)
d = s.shapes.add_shape(1, Inches(3.0), Inches(3.45), Inches(7.3), Inches(0.018))
d.fill.solid(); d.fill.fore_color.rgb = WHITE; d.line.fill.background()
txt(s, 'Presented by the Software Development Team',
    0.5, 3.62, 12.3, 0.45, size=15, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 'Africa Eye Foundation  ·  2025',
    0.5, 4.12, 12.3, 0.4, size=14, color=CYAN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS AEF HRM?
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
top_bar(s); bottom_bar(s)
title_bar(s, 'What is AEF HRM?', 'The problem it solves & what it delivers')
logo(s)

txt(s,
    'AEF HRM (LeaveDesk) is a fully web-based Human Resource Management System built '
    'specifically for Africa Eye Foundation. It replaces paper forms and manual processes '
    'with structured, automated, role-driven digital workflows.',
    0.5, 1.2, 12.3, 0.85, size=14, color=BLACK)

divider(s, 2.1)

# 4 stat boxes
stats = [('8+ Roles', 'Full role-based\naccess control'),
         ('6-Step Chain', 'Appraisal approval\nhierarchy'),
         ('100% Paperless', 'Leave, contracts\n& appraisals'),
         ('Auto PDF', 'Letters & reports\ngenerated instantly')]
for i, (val, lbl) in enumerate(stats):
    x = 0.4 + i * 3.1
    box(s, x, 2.2, 2.75, 1.55, NAVY)
    txt(s, val,  x, 2.25, 2.75, 0.75, size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    txt(s, lbl,  x, 2.95, 2.75, 0.75, size=12, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, 'Key Modules:', 0.5, 3.9, 12.0, 0.38, size=15, bold=True, color=NAVY)
modules = [
    '👤 Employee Management — profiles, documents, history',
    '🏖️ Leave Management — multi-role approval chain, auto balance, PDF letters',
    '📋 Personnel Appraisal — 6-step graded workflow with score overrides & PDF',
    '📄 Contracts & ⚖️ Discipline — issue, renew, sanctions linked to appraisal',
]
bullet(s, modules, t=4.3, size=13.5, gap=0.72)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LEAVE MANAGEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
top_bar(s); bottom_bar(s)
title_bar(s, 'Leave Management', 'Multi-step approval chain with auto PDF generation')
logo(s)

# Approval chain
steps = ['Employee\nSubmits', 'Supervisor', 'Unit Head', 'HR Manager', 'Director', '✓ Approved']
for i, step in enumerate(steps):
    x = 0.4 + i * 2.1
    bg = CYAN if i == 5 else NAVY
    box(s, x, 1.18, 1.82, 0.8, bg)
    txt(s, step, x, 1.18, 1.82, 0.8, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, '→', x + 1.82, 1.42, 0.28, 0.35, size=16, bold=True,
            color=NAVY, align=PP_ALIGN.CENTER)

divider(s, 2.1)

feats = [
    ('Leave Types', 'Annual, Sick, Maternity, Compassionate, Unpaid, WACS Residency & more'),
    ('Auto Balance', 'Leave balances tracked and deducted automatically — no manual counting'),
    ('PDF Letters', 'Professional approval letters generated instantly with all signatures'),
    ('Notifications', 'Email + in-app alerts sent to everyone involved at every workflow step'),
    ('Calendar View', 'Visual tracker showing who is on leave across the whole organisation'),
    ('Reject & Resubmit', 'Rejections notify employee with reason; employee can resubmit'),
]
for i, (lbl, desc) in enumerate(feats):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.45; y = 2.22 + row * 1.6
    box(s, x, y, 6.0, 1.4, WHITE)
    box(s, x, y, 1.6, 0.38, CYAN)
    txt(s, lbl, x, y, 1.6, 0.38, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, x + 0.15, y + 0.45, 5.65, 0.85, size=12.5, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — APPRAISAL MODULE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
top_bar(s); bottom_bar(s)
title_bar(s, 'Personnel Appraisal Module', 'End-to-end grading, overrides, PDF — fully digital')
logo(s)

# 6-step chain
chain = [('1', 'Employee'), ('2', 'Co-Worker'), ('3', 'Supervisor'),
         ('4', 'HR'), ('5', 'Director'), ('6', 'CEO')]
for i, (num, role) in enumerate(chain):
    x = 0.38 + i * 2.1
    bg = CYAN if i == 5 else NAVY
    box(s, x, 1.18, 1.9, 0.6, bg)
    txt(s, f'{num}. {role}', x, 1.18, 1.9, 0.6,
        size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(s, '▶', x + 1.9, 1.36, 0.2, 0.28, size=11,
            color=CYAN, align=PP_ALIGN.CENTER)

divider(s, 1.9)

# Scoring
txt(s, 'Scoring System', 0.5, 1.98, 12.0, 0.4, size=15, bold=True, color=NAVY)
scores = [
    ('Performance Factors',  '4 factors  ×  1–5', '12.5 pts max'),
    ('Attitude & Aptitude',  '6 factors  ×  1–5', '7.5 pts max'),
    ('Discipline Sanctions', 'Linked from HR records', '–1 per sanction'),
    ('Awards & Bonus',       'Employee of Month etc.', '+1 per award'),
]
for i, (cat, detail, pts) in enumerate(scores):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.45; y = 2.45 + row * 1.1
    box(s, x, y, 6.0, 0.95, WHITE)
    txt(s, cat, x + 0.15, y + 0.07, 2.8, 0.42, bold=True, size=13, color=NAVY)
    txt(s, detail, x + 0.15, y + 0.52, 2.8, 0.4, size=12, color=BLACK)
    txt(s, pts, x + 3.5, y + 0.22, 2.3, 0.45,
        size=13, bold=True, color=CYAN, align=PP_ALIGN.RIGHT)

divider(s, 4.72)

txt(s, 'Total = Performance + Attitude + Discipline + Awards  =  20 pts maximum',
    0.5, 4.79, 12.3, 0.45, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

txt(s, 'Score Override:', 0.5, 5.35, 2.5, 0.38, size=13, bold=True, color=NAVY)
txt(s,
    'HR, Admin Director, or CEO can modify the supervisor\'s grades. '
    'PDF shows both the original AND the override score with who changed it. '
    'Highest hierarchy is the final mark.',
    3.0, 5.35, 9.8, 0.7, size=13, color=BLACK)

txt(s, 'What the Supervisor sees:', 0.5, 6.15, 3.5, 0.38, size=13, bold=True, color=NAVY)
txt(s,
    'Employee self-assessment  →  Co-worker comment  →  Then fills appraiser rating tables '
    '+ live Total Ratings calculator  →  Signs.',
    4.0, 6.15, 8.8, 0.7, size=13, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNOLOGY & DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
top_bar(s); bottom_bar(s)
title_bar(s, 'Technology & Deployment', 'Built to be fast, secure, and always available')
logo(s)

tech_cards = [
    ('Backend', 'Python 3.11 · Django 4.x\nPostgreSQL · Gunicorn\nReportLab (PDF) · Pillow (signatures)'),
    ('Frontend', 'Bootstrap 5 · Bootstrap Icons\nSignaturePad.js (digital sigs)\nVanilla JS live score calculator'),
    ('Deployment', 'Railway cloud PaaS\nAuto-deploy from GitHub\nHTTPS enforced, zero downtime'),
    ('Email & Security', 'Resend / Anymail — transactional email\nRole-based access on every view\nCSRF, password hashing, env secrets'),
]
for i, (head, body) in enumerate(tech_cards):
    col = i % 2; row = i // 2
    x = 0.4 + col * 6.45; y = 1.18 + row * 2.55
    card(s, x, y, 6.0, 2.35, head, body)

divider(s, 5.38)
txt(s, 'GitHub push  →  Railway auto-deploys in under 2 minutes  ·  '
       'Accessible from any browser, any device, anywhere',
    0.5, 5.45, 12.3, 0.45, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

infra_pts = [
    'WhiteNoise serves static files — no CDN required',
    'Signatures stored as base64 in DB — no file storage dependency in production',
    'All secrets in Railway environment variables — nothing hardcoded in source',
]
bullet(s, infra_pts, t=5.98, size=13.5, gap=0.44)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE DEMO & NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.33, 7.5, NAVY)
box(s, 0, 0, 13.33, 0.07, CYAN)
box(s, 0, 7.43, 13.33, 0.07, CYAN)
logo(s, l=0.4, t=0.2, w=1.9, h=0.68)
bottom_bar(s)

txt(s, 'Live Demo  &  Next Steps', 2.2, 0.12, 10.5, 0.65,
    size=30, bold=True, color=WHITE)
d = s.shapes.add_shape(1, Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.018))
d.fill.solid(); d.fill.fore_color.rgb = CYAN; d.line.fill.background()

# Demo column
box(s, 0.4, 1.02, 5.9, 0.48, CYAN)
txt(s, '▶  Live Demo Walkthrough (5 min)', 0.55, 1.02, 5.75, 0.48,
    size=14, bold=True, color=NAVY)
demo_steps = [
    'Login as Employee — submit a leave request',
    'Login as Supervisor — approve, see dashboard pending',
    'Open Appraisal — employee fill → co-worker → supervisor grades',
    'Show live Total Ratings calculator updating in real time',
    'Download the appraisal PDF — see all signatures & scores',
    'Show contract issuance & notification email',
]
for i, step in enumerate(demo_steps):
    txt(s, f'  {i+1}.  {step}', 0.5, 1.6 + i * 0.85, 5.7, 0.78,
        size=13, color=WHITE)

# Next steps column
box(s, 6.9, 1.02, 5.9, 0.48, CYAN)
txt(s, '📌  Immediate Next Steps', 7.05, 1.02, 5.75, 0.48,
    size=14, bold=True, color=NAVY)
next_steps = [
    'Apply DB migration on Railway',
    'Onboard all staff with correct roles',
    'Configure AEF email domain in Resend',
    'HR initiates Trimester 1 – 2025 appraisal cycle',
    'Collect digital signatures from chain',
    'HR distributes final appraisal PDFs',
]
for i, step in enumerate(next_steps):
    c = CYAN if i < 2 else WHITE
    txt(s, f'  {"✅" if i < 1 else "📌"}  {step}', 6.9, 1.6 + i * 0.85, 5.7, 0.78,
        size=13, color=c)

d2 = s.shapes.add_shape(1, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.018))
d2.fill.solid(); d2.fill.fore_color.rgb = CYAN; d2.line.fill.background()
txt(s, 'Thank you  —  Questions are Welcome',
    0.4, 6.78, 12.5, 0.45, size=18, bold=True,
    color=CYAN, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), 'AEF_HRM_Presentation.pptx')
prs.save(out)
print(f'Saved: {out}')
